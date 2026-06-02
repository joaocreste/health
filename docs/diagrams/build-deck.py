#!/usr/bin/env python3
"""Rebuild Lumen-Health-Architecture.pptx from the SVGs in this directory.

Cover is full-bleed; each diagram is rasterized at high resolution and
fit-to-slide, centered, on a matching #0A1428 background. 16:9.
Run from docs/diagrams/:  python3 build-deck.py
"""
import subprocess, os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
BG = RGBColor(0x0A, 0x14, 0x28)
RENDER_W = 2600  # px width for rasterization

# (svg file, full_bleed?)  — order defines slide order
SLIDES = [
    ("00-cover.svg", True),
    ("01-system-architecture.svg", False),
    ("02-ingestion-etl-pipeline.svg", False),
    ("03-data-treatment-routing.svg", False),
    ("04-data-residency-locations.svg", False),
    ("05-ai-llm-architecture.svg", False),
    ("06-database-schema-er.svg", False),
    ("07-auth-rbac.svg", False),
    ("08-multitenancy-rendering.svg", False),
    ("09-compliance-state-machine.svg", False),
    ("10-storage-file-types.svg", False),
    ("11-neon-table-structures.svg", False),
]

def rasterize(svg):
    png = os.path.join(HERE, "_build_" + svg.replace(".svg", ".png"))
    subprocess.run(["rsvg-convert", "-w", str(RENDER_W),
                    os.path.join(HERE, svg), "-o", png], check=True)
    return png

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

for svg, full_bleed in SLIDES:
    png = rasterize(svg)
    slide = prs.slides.add_slide(blank)
    # dark background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    iw, ih = Image.open(png).size
    if full_bleed:
        # cover fills the slide (it is authored 16:9)
        slide.shapes.add_picture(png, 0, 0, width=SW, height=SH)
    else:
        # fit-to-slide, centered, preserving aspect
        scale = min(SW / iw, SH / ih)
        w, h = int(iw * scale), int(ih * scale)
        left, top = int((SW - w) / 2), int((SH - h) / 2)
        slide.shapes.add_picture(png, Emu(left), Emu(top), width=Emu(w), height=Emu(h))

out = os.path.join(HERE, "Lumen-Health-Architecture.pptx")
prs.save(out)
# clean intermediate pngs
for svg, _ in SLIDES:
    p = os.path.join(HERE, "_build_" + svg.replace(".svg", ".png"))
    if os.path.exists(p):
        os.remove(p)
print(f"Wrote {out} — {len(SLIDES)} slides")
