#!/usr/bin/env python3
"""
Builds infrastructure.pptx — a 12-slide deck mirroring infrastructure.html.

Each diagram element is a native PowerPoint shape so individual rectangles,
arrows and labels can be lifted into a pitch deck. Brand-consistent with the
Lumen Health webapp (navy / slate / gold / steel / off-white).

Run from the project root:
    python3 bin/build-infrastructure-pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
from pathlib import Path

# ─── Brand tokens ────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)
SLATE       = RGBColor(0x1E, 0x2D, 0x3D)
GOLD        = RGBColor(0xB8, 0x95, 0x4A)
GOLD_SOFT   = RGBColor(0xFF, 0xF6, 0xE5)
GOLD_LINE   = RGBColor(0xE0, 0xC6, 0x81)
STEEL       = RGBColor(0x7A, 0x8F, 0xA6)
OFFWHITE    = RGBColor(0xF9, 0xF7, 0xF4)
SURFACE     = RGBColor(0xF4, 0xF7, 0xFA)
BORDER      = RGBColor(0xE5, 0xE2, 0xDC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1A, 0x21, 0x29)
INK_SOFT    = RGBColor(0x6E, 0x7B, 0x8A)
BLUE_50     = RGBColor(0xEE, 0xF5, 0xFA)
BLUE_100    = RGBColor(0xD8, 0xE8, 0xF2)
BLUE_500    = RGBColor(0x3E, 0x7C, 0xA3)
BLUE_700    = RGBColor(0x24, 0x4E, 0x6E)
BLUE_800    = RGBColor(0x1B, 0x3B, 0x54)

FONT_DISPLAY = "Raleway"           # falls back gracefully if absent
FONT_BODY    = "IBM Plex Sans"
FONT_MONO    = "IBM Plex Mono"

TOTAL_SLIDES = 12

# ─── Helpers ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, *, fill=WHITE, stroke=NAVY, stroke_w=1.0,
             rounded=False, dash=False):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_kind, x, y, w, h)
    if rounded:
        shp.adjustments[0] = 0.06
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(stroke_w)
        if dash:
            shp.line.dash_style = 7  # MSO_LINE.DASH
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=11, bold=False,
             color=SLATE, align="left", anchor="top", italic=False,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_arrow(slide, x1, y1, x2, y2, *, color=SLATE, weight=1.5, dash=False):
    """Straight connector with arrowhead at the end."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        conn.line.dash_style = 7
    # Add arrowhead on the tail (line goes from x1,y1 to x2,y2; arrow at x2,y2)
    line = conn.line._get_or_add_ln()
    tail = etree.SubElement(line, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn

def add_line(slide, x1, y1, x2, y2, *, color=BORDER, weight=1.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        conn.line.dash_style = 7
    return conn

# ─── Slide scaffold ──────────────────────────────────────────────

def slide_header(slide, idx, eyebrow_label, title_text, subtitle_text=None):
    """Standard header: gold eyebrow top-left, page counter top-right, big title."""
    # Top gold accent strip
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.07),
             fill=GOLD, stroke=None)

    # Eyebrow (top-left)
    add_text(slide, Inches(0.5), Inches(0.28), Inches(8), Inches(0.3),
             eyebrow_label, font=FONT_MONO, size=10, bold=True,
             color=GOLD)

    # Page counter (top-right)
    add_text(slide, Inches(11.5), Inches(0.28), Inches(1.3), Inches(0.3),
             f"{idx:02d} / {TOTAL_SLIDES:02d}",
             font=FONT_MONO, size=10, bold=False, color=STEEL, align="right")

    # Title
    add_text(slide, Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.85),
             title_text, font=FONT_DISPLAY, size=30, bold=True,
             color=NAVY)

    if subtitle_text:
        add_text(slide, Inches(0.5), Inches(1.42), Inches(12.3), Inches(0.5),
                 subtitle_text, font=FONT_BODY, size=14, bold=False,
                 color=SLATE)

def slide_footer(slide, footer_text="Lumen Health · Infrastructure"):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             footer_text, font=FONT_MONO, size=9, color=INK_SOFT)

# ─── Build the presentation ─────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════════
# Slide 1 · TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Full navy background
add_rect(s, 0, 0, prs.slide_width, prs.slide_height,
         fill=NAVY, stroke=None)

# Gold top accent
add_rect(s, 0, Inches(3.3), Inches(13.333), Inches(0.08),
         fill=GOLD, stroke=None)

# Eyebrow
add_text(s, Inches(1), Inches(2.8), Inches(11), Inches(0.4),
         "LUMEN HEALTH · ENGINEERING DOCS",
         font=FONT_MONO, size=14, bold=True, color=GOLD)

# Big title
add_text(s, Inches(1), Inches(3.5), Inches(11), Inches(1.4),
         "Infrastructure & ETL workflow",
         font=FONT_DISPLAY, size=54, bold=True, color=WHITE)

# Subtitle
add_text(s, Inches(1), Inches(4.9), Inches(11), Inches(0.9),
         "How raw wearable exports, lab PDFs and DICOM bundles become charts, AI insights and bilingual narrative pages — end to end.",
         font=FONT_BODY, size=18, color=RGBColor(0xCF, 0xD7, 0xE0))

# Footer chips
add_text(s, Inches(1), Inches(6.4), Inches(11), Inches(0.4),
         "Cloudflare Pages + Workers   ·   Neon Postgres (Frankfurt)   ·   Cloudflare R2   ·   Anthropic Claude",
         font=FONT_MONO, size=12, color=GOLD)


# ════════════════════════════════════════════════════════════════
# Slide 2 · AT A GLANCE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 2, "01 · AT A GLANCE",
             "What lives where",
             "A build-time path runs on the engineer's laptop; runtime is fully Cloudflare-native serverless.")

# 4 stat cards in one row
cards = [
    ("Data sources", "7", "device + manual streams",
     "Oura · Apple · Withings · CGM · DICOM · CSVs · writings"),
    ("Database tables", "31", "Postgres tables · 10 domains",
     "Drizzle ORM · 9 enums · schema in db/schema.ts"),
    ("Claude models", "3", "tiers in production",
     "Opus 4.7 (chat) · Sonnet 4.6 (ingest) · Haiku 4.5 (classify)"),
    ("Deploy cadence", "Manual", "wrangler CLI",
     "Git push ≠ deploy. Live URL updates after wrangler pages deploy"),
]

card_y = Inches(2.4)
card_h = Inches(2.6)
card_w = Inches(2.95)
card_gap = Inches(0.18)
card_x0 = Inches(0.5)

for i, (label, value, sub, foot) in enumerate(cards):
    x = card_x0 + (card_w + card_gap) * i
    # Gold left bar
    add_rect(s, x, card_y, Inches(0.08), card_h, fill=GOLD, stroke=None)
    # Card background
    add_rect(s, x + Inches(0.08), card_y, card_w - Inches(0.08), card_h,
             fill=WHITE, stroke=BORDER)
    # Label
    add_text(s, x + Inches(0.3), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3),
             label.upper(), font=FONT_MONO, size=10, bold=True, color=STEEL)
    # Value
    add_text(s, x + Inches(0.3), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.8),
             value, font=FONT_DISPLAY, size=36, bold=True, color=NAVY)
    # Sub
    add_text(s, x + Inches(0.3), card_y + Inches(1.4), card_w - Inches(0.4), Inches(0.3),
             sub, font=FONT_BODY, size=11, color=STEEL)
    # Foot
    add_text(s, x + Inches(0.3), card_y + Inches(1.75), card_w - Inches(0.4), Inches(0.8),
             foot, font=FONT_BODY, size=11.5, color=SLATE)

# Insight strip below
add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0),
         fill=GOLD_SOFT, stroke=GOLD_LINE)
add_rect(s, Inches(0.5), Inches(5.5), Inches(0.06), Inches(1.0),
         fill=GOLD, stroke=None)
add_text(s, Inches(0.75), Inches(5.65), Inches(12), Inches(0.3),
         "DESIGN PRINCIPLE", font=FONT_MONO, size=10, bold=True, color=GOLD)
add_text(s, Inches(0.75), Inches(5.95), Inches(12), Inches(0.5),
         "Separate the build path from the runtime path. Static charts ship as files. Only AI, auth and data mutations touch the Worker.",
         font=FONT_BODY, size=13, color=SLATE)

slide_footer(s)


# ════════════════════════════════════════════════════════════════
# Slide 3 · TOPOLOGY DIAGRAM (build-time + runtime overview)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 3, "02 · BIG PICTURE",
             "End-to-end topology",
             "Build-time (left) converts raw exports into static artifacts. Runtime (right) serves them and dispatches API calls.")

# Two-zone label band
add_text(s, Inches(0.5), Inches(2.05), Inches(6), Inches(0.25),
         "BUILD-TIME · LOCAL LAPTOP", font=FONT_MONO, size=9.5, bold=True, color=STEEL)
add_line(s, Inches(0.5), Inches(2.28), Inches(6.5), Inches(2.28), color=BORDER)
add_text(s, Inches(6.7), Inches(2.05), Inches(6), Inches(0.25),
         "RUNTIME · CLOUDFLARE-NATIVE", font=FONT_MONO, size=9.5, bold=True, color=STEEL)
add_line(s, Inches(6.7), Inches(2.28), Inches(12.85), Inches(2.28), color=BORDER)

# ── Build-time row ─────────────────────────────────────────────
# Data sources stack (col 1)
sources = ["Oura Ring · CSVs", "Apple Health · ECG + XML", "Withings · scale + BP",
           "FreeStyle Libre CGM", "DICOM workstation", "Manual CSVs + writings"]
src_y = Inches(2.5)
for i, name in enumerate(sources):
    stroke = NAVY if i < 4 else STEEL
    dash = i >= 4
    y = src_y + Inches(0.42) * i
    add_rect(s, Inches(0.5), y, Inches(2.5), Inches(0.35),
             fill=WHITE, stroke=stroke, dash=dash)
    add_text(s, Inches(0.65), y + Inches(0.04), Inches(2.3), Inches(0.3),
             name, font=FONT_BODY, size=10, bold=True, color=NAVY, anchor="middle")

# Arrow → ETL
add_arrow(s, Inches(3.05), Inches(3.5), Inches(3.4), Inches(3.5), color=SLATE)

# ETL box (col 2)
add_rect(s, Inches(3.45), Inches(2.5), Inches(2.4), Inches(2.6),
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.6)
add_text(s, Inches(3.45), Inches(2.6), Inches(2.4), Inches(0.3),
         "bin/extract.py", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(3.45), Inches(2.92), Inches(2.4), Inches(0.25),
         "Python 3 · CLI · stdlib", font=FONT_BODY, size=10, color=SLATE, align="center")
add_line(s, Inches(3.65), Inches(3.3), Inches(5.65), Inches(3.3), color=GOLD_LINE)
bullets = [
    "reads patient folder",
    "aggregates per-day series",
    "Tukey boxplots, percentiles",
    "ECG classification",
    "rolling means / IQRs",
]
for i, b in enumerate(bullets):
    add_text(s, Inches(3.55), Inches(3.4) + Inches(0.27) * i, Inches(2.3), Inches(0.25),
             "• " + b, font=FONT_BODY, size=10, color=SLATE)

# Arrow → artifacts
add_arrow(s, Inches(5.9), Inches(3.5), Inches(6.25), Inches(3.5), color=SLATE)

# Artifacts column (col 3)
arts = [
    ("web/assets/data.js",  "daily time-series"),
    ("web/assets/metrics.json", "summary stats"),
    ("web/*.html (EN + PT)", "static pages"),
    ("web/scans/*.png", "DICOM slices"),
    ("web/_worker.js", "Worker bundle"),
]
art_y = Inches(2.5)
for i, (name, sub) in enumerate(arts):
    y = art_y + Inches(0.5) * i
    add_rect(s, Inches(6.3), y, Inches(2.5), Inches(0.45),
             fill=WHITE, stroke=NAVY)
    add_text(s, Inches(6.45), y + Inches(0.04), Inches(2.3), Inches(0.22),
             name, font=FONT_MONO, size=9.5, bold=True, color=NAVY)
    add_text(s, Inches(6.45), y + Inches(0.25), Inches(2.3), Inches(0.18),
             sub, font=FONT_BODY, size=9.5, color=STEEL)

# Arrow → deploy
add_arrow(s, Inches(8.85), Inches(3.5), Inches(9.25), Inches(3.5), color=GOLD, weight=1.8)
add_text(s, Inches(8.85), Inches(3.2), Inches(0.5), Inches(0.25),
         "DEPLOY", font=FONT_MONO, size=8, bold=True, color=GOLD, align="center")

# Deploy + Pages box
add_rect(s, Inches(9.3), Inches(2.5), Inches(3.55), Inches(2.6),
         fill=NAVY, stroke=NAVY)
add_text(s, Inches(9.3), Inches(2.65), Inches(3.55), Inches(0.3),
         "wrangler CLI", font=FONT_DISPLAY, size=14, bold=True, color=WHITE, align="center")
add_text(s, Inches(9.3), Inches(2.95), Inches(3.55), Inches(0.25),
         "npx wrangler pages deploy web",
         font=FONT_MONO, size=9.5, color=RGBColor(0xCF, 0xD7, 0xE0), align="center")
add_line(s, Inches(9.5), Inches(3.35), Inches(12.65), Inches(3.35),
         color=RGBColor(0x4F, 0x6A, 0x80))
add_text(s, Inches(9.3), Inches(3.45), Inches(3.55), Inches(0.35),
         "Cloudflare Pages", font=FONT_DISPLAY, size=14, bold=True, color=WHITE, align="center")
add_text(s, Inches(9.3), Inches(3.78), Inches(3.55), Inches(0.25),
         "jc-advisory-health.pages.dev",
         font=FONT_MONO, size=9.5, color=GOLD, align="center")
deploy_notes = [
    "• serves static HTML / JS / PNG",
    "• routes /api/* → Worker",
    "• global edge CDN",
]
for i, n in enumerate(deploy_notes):
    add_text(s, Inches(9.45), Inches(4.15) + Inches(0.24) * i, Inches(3.4), Inches(0.22),
             n, font=FONT_BODY, size=9.5,
             color=RGBColor(0xCF, 0xD7, 0xE0))

# ── Runtime row ────────────────────────────────────────────────
add_line(s, Inches(0.5), Inches(5.35), Inches(12.85), Inches(5.35),
         color=BORDER, dash=True)
add_text(s, Inches(0.5), Inches(5.45), Inches(8), Inches(0.25),
         "RUNTIME · USER-FACING REQUEST PATH",
         font=FONT_MONO, size=9.5, bold=True, color=STEEL)

# User
add_rect(s, Inches(0.5), Inches(5.85), Inches(1.8), Inches(0.8),
         fill=WHITE, stroke=NAVY)
add_text(s, Inches(0.5), Inches(5.95), Inches(1.8), Inches(0.3),
         "User browser", font=FONT_DISPLAY, size=11, bold=True, color=NAVY, align="center")
add_text(s, Inches(0.5), Inches(6.22), Inches(1.8), Inches(0.4),
         "role · EN/PT", font=FONT_BODY, size=10, color=SLATE, align="center")

add_arrow(s, Inches(2.35), Inches(6.25), Inches(2.7), Inches(6.25), color=SLATE)

# CF Pages
add_rect(s, Inches(2.8), Inches(5.85), Inches(2.5), Inches(0.8),
         fill=BLUE_50, stroke=NAVY)
add_text(s, Inches(2.8), Inches(5.95), Inches(2.5), Inches(0.3),
         "Cloudflare Pages", font=FONT_DISPLAY, size=11, bold=True, color=NAVY, align="center")
add_text(s, Inches(2.8), Inches(6.22), Inches(2.5), Inches(0.4),
         "edge CDN · static", font=FONT_BODY, size=10, color=SLATE, align="center")

add_arrow(s, Inches(5.35), Inches(6.25), Inches(5.75), Inches(6.25), color=GOLD, weight=1.6)
add_text(s, Inches(5.3), Inches(5.9), Inches(0.5), Inches(0.22),
         "/api/*", font=FONT_MONO, size=8.5, bold=True, color=GOLD, align="center")

# Worker
add_rect(s, Inches(5.85), Inches(5.85), Inches(2.3), Inches(0.8),
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.5)
add_text(s, Inches(5.85), Inches(5.95), Inches(2.3), Inches(0.3),
         "_worker.js", font=FONT_DISPLAY, size=11, bold=True, color=NAVY, align="center")
add_text(s, Inches(5.85), Inches(6.22), Inches(2.3), Inches(0.4),
         "router · auth · fan-out", font=FONT_BODY, size=10, color=SLATE, align="center")

# Fan out to services
add_arrow(s, Inches(8.2), Inches(5.95), Inches(8.85), Inches(5.95), color=SLATE)
add_arrow(s, Inches(8.2), Inches(6.25), Inches(8.85), Inches(6.25), color=SLATE)
add_arrow(s, Inches(8.2), Inches(6.55), Inches(8.85), Inches(6.55), color=SLATE)

# Three services
svc_y = Inches(5.78)
for i, (name, sub) in enumerate([
    ("Neon Postgres · EU", "32 tables · Drizzle ORM"),
    ("Cloudflare R2", "blobs · binding R2_BUCKET"),
    ("Anthropic Claude", "Opus · Sonnet · Haiku"),
]):
    y = svc_y + Inches(0.31) * i
    add_rect(s, Inches(8.9), y, Inches(3.95), Inches(0.28),
             fill=BLUE_50, stroke=NAVY)
    add_text(s, Inches(9.0), y + Inches(0.03), Inches(2.4), Inches(0.22),
             name, font=FONT_BODY, size=9.5, bold=True, color=NAVY, anchor="middle")
    add_text(s, Inches(11.0), y + Inches(0.03), Inches(1.8), Inches(0.22),
             sub, font=FONT_MONO, size=8.5, color=BLUE_700, anchor="middle", align="right")

slide_footer(s)


# ════════════════════════════════════════════════════════════════
# Slide 4 · DATA SOURCES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 4, "03 · DATA SOURCES",
             "Where the raw signal comes from",
             "All inputs live under Patients/<Patient>/ at the repo root. The script accepts JC_DATA_DIR to point elsewhere.")

# Table of sources
hdr = ["Source", "Folder", "Format", "Refreshed", "What's extracted"]
rows = [
    ("Oura Ring",       "Oura/App Data/",                      "~45 CSVs",      "monthly", "Sleep stages · HRV · RHR · steps · stress · resilience · workouts · tags"),
    ("Apple Health",    "Apple Health/",                       "XML + ECG + GPX","monthly", "63 ECGs · 188 GPX routes · 866 MB XML kept local"),
    ("Withings",        "Withings/",                           "~20 CSVs",      "monthly", "121 bio-impedance readings · 249 BP cuff readings since Nov 2025"),
    ("FreeStyle Libre", "glucose_timeline.csv",                "CSV · 5-min",   "per sensor", "1,712 continuous readings, 26 Apr → 2 May 2026"),
    ("DICOM",           "Imagery/",                            "DICOM stacks",  "per visit", "Pre-rendered to PNG + manifest.json before deploy"),
    ("Lab PDFs",        "Blood/",                              "PDF",           "per visit", "Transcribed to Postgres via the ingest API (Claude Sonnet)"),
    ("Manual CSVs",     "root of patient folder",              "small CSVs",    "manual edit", "medications · supplements · AUDIT · clinical_history · wheel_of_life"),
    ("Writings",        "Writtings/ · Therapy Transcripts/",   ".txt · .json",  "as authored", "Patient first-person material — quoted verbatim, never auto-rewritten"),
]

# Header band
hdr_y = Inches(2.2)
col_w = [Inches(1.6), Inches(2.3), Inches(1.6), Inches(1.5), Inches(5.3)]
col_x = [Inches(0.5)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.35),
         fill=NAVY, stroke=None)
for i, label in enumerate(hdr):
    add_text(s, col_x[i] + Inches(0.12), hdr_y + Inches(0.05),
             col_w[i] - Inches(0.2), Inches(0.25),
             label.upper(), font=FONT_MONO, size=9.5, bold=True, color=GOLD, anchor="middle")

# Rows
row_y = hdr_y + Inches(0.35)
row_h = Inches(0.5)
for ri, row in enumerate(rows):
    if ri % 2 == 1:
        add_rect(s, Inches(0.5), row_y, Inches(12.3), row_h,
                 fill=SURFACE, stroke=None)
    for ci, val in enumerate(row):
        font = FONT_MONO if ci in (1, 2) else FONT_BODY
        bold = ci == 0
        size = 10.5 if ci != 4 else 10
        color = NAVY if ci == 0 else SLATE
        if ci == 1:
            color = BLUE_700
        add_text(s, col_x[ci] + Inches(0.12), row_y + Inches(0.07),
                 col_w[ci] - Inches(0.2), row_h - Inches(0.14),
                 val, font=font, size=size, bold=bold, color=color,
                 anchor="middle", line_spacing=1.1)
    row_y += row_h
    # row separator
    add_line(s, Inches(0.5), row_y, Inches(12.8), row_y, color=BORDER)

slide_footer(s)


# ════════════════════════════════════════════════════════════════
# Slide 5 · ETL PIPELINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 5, "04 · BUILD-TIME ETL",
             "Python pipeline — bin/extract.py",
             "Reads every CSV / JSON / PDF under the patient folder, computes per-day aggregates, writes two artifacts back into web/assets/.")

# 3 columns: inputs, transforms, outputs
# Inputs
add_text(s, Inches(0.5), Inches(2.1), Inches(4), Inches(0.25),
         "INPUTS · Patients/Joao Victor Creste/",
         font=FONT_MONO, size=9.5, bold=True, color=STEEL)
add_line(s, Inches(0.5), Inches(2.33), Inches(4.3), Inches(2.33), color=BORDER)

inputs = [
    "Oura/App Data/sleepmodel.csv",
    "Oura/App Data/dailyactivity.csv",
    "Oura/App Data/dailysleep.csv",
    "Oura/App Data/dailystress.csv",
    "Oura/App Data/heartrate.csv",
    "Withings/weight.csv",
    "Withings/bp.csv",
    "Apple Health/electrocardiograms/",
    "glucose_timeline.csv",
    "medications.csv · audit.csv · …",
]
for i, name in enumerate(inputs):
    y = Inches(2.55) + Inches(0.34) * i
    stroke = STEEL if i >= 8 else NAVY
    add_rect(s, Inches(0.5), y, Inches(4), Inches(0.28),
             fill=WHITE, stroke=stroke, dash=(i >= 8))
    add_text(s, Inches(0.6), y + Inches(0.04), Inches(3.9), Inches(0.2),
             name, font=FONT_MONO, size=9, color=BLUE_700, anchor="middle")

# Arrow → transforms
add_arrow(s, Inches(4.55), Inches(4.1), Inches(4.95), Inches(4.1), color=SLATE)

# Transforms box
add_text(s, Inches(5), Inches(2.1), Inches(4), Inches(0.25),
         "TRANSFORMS",
         font=FONT_MONO, size=9.5, bold=True, color=STEEL)
add_line(s, Inches(5), Inches(2.33), Inches(9.05), Inches(2.33), color=BORDER)

add_rect(s, Inches(5), Inches(2.55), Inches(4), Inches(4.4),
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.6)
add_text(s, Inches(5), Inches(2.7), Inches(4), Inches(0.35),
         "bin/extract.py", font=FONT_DISPLAY, size=18, bold=True, color=NAVY, align="center")
add_text(s, Inches(5), Inches(3.05), Inches(4), Inches(0.25),
         "Python 3 · pure stdlib", font=FONT_BODY, size=11, color=SLATE, align="center")
add_line(s, Inches(5.2), Inches(3.45), Inches(8.8), Inches(3.45), color=GOLD_LINE)

transforms = [
    ("PER-NIGHT",   ["Tukey boxplots for sleep stages",
                     "HRV / RHR monthly aggregates"]),
    ("PER-DAY",     ["steps · sedentary · active cal",
                     "stress min · resilience score"]),
    ("PER-WEEK / MONTH", ["BP weekly bands (±1 SD, ±2 SD)",
                          "HR-by-time-of-day folds"]),
    ("PARSE",       ["ECG classification + dates",
                     "Oura self-tags (alcohol, valium)"]),
]
y = Inches(3.6)
for heading, bullets in transforms:
    add_text(s, Inches(5.15), y, Inches(3.7), Inches(0.22),
             heading, font=FONT_MONO, size=8.5, bold=True, color=GOLD)
    y += Inches(0.22)
    for b in bullets:
        add_text(s, Inches(5.15), y, Inches(3.7), Inches(0.22),
                 b, font=FONT_BODY, size=10, color=SLATE)
        y += Inches(0.22)
    y += Inches(0.04)

# Arrow → outputs
add_arrow(s, Inches(9.1), Inches(4.1), Inches(9.5), Inches(4.1), color=GOLD, weight=1.8)

# Outputs
add_text(s, Inches(9.55), Inches(2.1), Inches(3.3), Inches(0.25),
         "OUTPUTS · web/assets/",
         font=FONT_MONO, size=9.5, bold=True, color=STEEL)
add_line(s, Inches(9.55), Inches(2.33), Inches(12.85), Inches(2.33), color=BORDER)

outputs = [
    ("data.js",
     "~74 KB · WEIGHT, HRV_RHR,",
     "STEPS, SLEEP_BOX, STRESS_RES …",
     "/* generated YYYY-MM-DD */"),
    ("metrics.json",
     "~96 KB · medians · IQRs · n · ranges",
     "consumed by HTML inline scripts",
     '{ "sleep": {…}, "ecgs": {…} }'),
]
y = Inches(2.7)
for name, l1, l2, code in outputs:
    add_rect(s, Inches(9.55), y, Inches(3.3), Inches(1.4),
             fill=WHITE, stroke=NAVY)
    add_text(s, Inches(9.7), y + Inches(0.12), Inches(3), Inches(0.3),
             name, font=FONT_DISPLAY, size=14, bold=True, color=NAVY)
    add_text(s, Inches(9.7), y + Inches(0.48), Inches(3), Inches(0.22),
             l1, font=FONT_BODY, size=10, color=SLATE)
    add_text(s, Inches(9.7), y + Inches(0.7), Inches(3), Inches(0.22),
             l2, font=FONT_BODY, size=10, color=SLATE)
    add_text(s, Inches(9.7), y + Inches(0.97), Inches(3), Inches(0.22),
             code, font=FONT_MONO, size=9, color=BLUE_700)
    y += Inches(1.6)

# Bottom callout
add_rect(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         fill=NAVY, stroke=None)
add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.25),
         "PIPELINE IS PURE-STDLIB PYTHON · IDEMPOTENT · RUN MANUALLY · NO SCHEDULED JOBS · NO CLOUD COMPUTE",
         font=FONT_MONO, size=9, bold=True, color=GOLD, anchor="middle")


# ════════════════════════════════════════════════════════════════
# Slide 6 · RUNTIME REQUEST FLOW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 6, "05 · RUNTIME",
             "What happens when a user clicks",
             "Static pages stream from the edge cache. The Worker only wakes for /api/* — chat, dashboards, ingest, admin.")

# Top row: User → Pages → Worker → three services
y_main = Inches(2.6)
box_h = Inches(1.0)

# User
add_rect(s, Inches(0.5), y_main, Inches(2), box_h, fill=WHITE, stroke=NAVY)
add_text(s, Inches(0.5), y_main + Inches(0.15), Inches(2), Inches(0.3),
         "User browser", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(0.5), y_main + Inches(0.5), Inches(2), Inches(0.4),
         "role: admin · doctor · patient", font=FONT_BODY, size=10, color=SLATE, align="center")

add_arrow(s, Inches(2.55), y_main + Inches(0.5), Inches(2.95), y_main + Inches(0.5), color=SLATE)
add_text(s, Inches(2.55), y_main + Inches(0.18), Inches(0.5), Inches(0.2),
         "HTTPS", font=FONT_MONO, size=8.5, bold=True, color=STEEL, align="center")

# CF Pages
add_rect(s, Inches(3), y_main, Inches(2.5), box_h,
         fill=BLUE_50, stroke=NAVY)
add_text(s, Inches(3), y_main + Inches(0.15), Inches(2.5), Inches(0.3),
         "Cloudflare Pages", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(3), y_main + Inches(0.5), Inches(2.5), Inches(0.4),
         "edge CDN · static cache", font=FONT_BODY, size=10, color=SLATE, align="center")

add_arrow(s, Inches(5.55), y_main + Inches(0.5), Inches(5.95), y_main + Inches(0.5), color=GOLD, weight=1.8)
add_text(s, Inches(5.45), y_main + Inches(0.18), Inches(0.7), Inches(0.2),
         "/api/*", font=FONT_MONO, size=8.5, bold=True, color=GOLD, align="center")

# Worker
add_rect(s, Inches(6), y_main, Inches(2.5), box_h,
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.6)
add_text(s, Inches(6), y_main + Inches(0.15), Inches(2.5), Inches(0.3),
         "_worker.js", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(6), y_main + Inches(0.5), Inches(2.5), Inches(0.4),
         "~1,200 LOC · single file", font=FONT_BODY, size=10, color=SLATE, align="center")

# Fan-out arrows
svc_x = Inches(9.5)
svc_w = Inches(3.35)
svc_y = [y_main - Inches(0.5), y_main + Inches(0.35), y_main + Inches(1.2)]
for sy in svc_y:
    add_arrow(s, Inches(8.55), y_main + Inches(0.5), svc_x, sy + Inches(0.25),
              color=SLATE)

# 3 services
services = [
    ("Neon Postgres · Frankfurt", "32 tables · Drizzle ORM",
     "@neondatabase/serverless · env.DATABASE_URL"),
    ("Cloudflare R2", "binding R2_BUCKET",
     "blobs · per-patient prefixes · GDPR cascade"),
    ("Anthropic Claude API", "@anthropic-ai/sdk",
     "Opus 4.7 · Sonnet 4.6 · Haiku 4.5"),
]
for sy, (name, l2, l3) in zip(svc_y, services):
    add_rect(s, svc_x, sy, svc_w, Inches(0.85),
             fill=BLUE_50, stroke=NAVY)
    add_text(s, svc_x + Inches(0.15), sy + Inches(0.08), svc_w - Inches(0.2), Inches(0.25),
             name, font=FONT_DISPLAY, size=12, bold=True, color=NAVY)
    add_text(s, svc_x + Inches(0.15), sy + Inches(0.35), svc_w - Inches(0.2), Inches(0.22),
             l2, font=FONT_BODY, size=10, color=SLATE)
    add_text(s, svc_x + Inches(0.15), sy + Inches(0.58), svc_w - Inches(0.2), Inches(0.22),
             l3, font=FONT_MONO, size=9, color=BLUE_700)

# Endpoint strip
add_line(s, Inches(0.5), Inches(5.5), Inches(12.85), Inches(5.5),
         color=BORDER, dash=True)
add_text(s, Inches(0.5), Inches(5.6), Inches(8), Inches(0.25),
         "WORKER ENDPOINT CLUSTER",
         font=FONT_MONO, size=9.5, bold=True, color=STEEL)
endpoints = (
    "/api/chat   ·   /api/me   ·   /api/login   ·   /api/patient-summary   ·   /api/patient-exams\n"
    "/api/patient-dashboard   ·   /api/patient-dashboard-build   ·   /api/patient-wipe-data   ·   /api/patients   ·   /api/admin/*"
)
add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7),
         endpoints, font=FONT_MONO, size=11, color=BLUE_700, line_spacing=1.4)

# Bottom split callout
add_rect(s, Inches(0.5), Inches(6.85), Inches(6), Inches(0.55),
         fill=BLUE_50, stroke=BLUE_500, stroke_w=1.3)
add_text(s, Inches(0.65), Inches(6.9), Inches(5.8), Inches(0.25),
         "STATIC PATH · ~95% OF TRAFFIC",
         font=FONT_MONO, size=9, bold=True, color=BLUE_500)
add_text(s, Inches(0.65), Inches(7.13), Inches(5.8), Inches(0.25),
         "Edge cache hit · ~zero cost · &lt; 50 ms anywhere",
         font=FONT_BODY, size=10.5, color=SLATE)

add_rect(s, Inches(6.85), Inches(6.85), Inches(6), Inches(0.55),
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.3)
add_text(s, Inches(7), Inches(6.9), Inches(5.8), Inches(0.25),
         "WORKER PATH · DATA + AI",
         font=FONT_MONO, size=9, bold=True, color=GOLD)
add_text(s, Inches(7), Inches(7.13), Inches(5.8), Inches(0.25),
         "V8 isolate · ~5 ms cold start · 30 s CPU cap",
         font=FONT_BODY, size=10.5, color=SLATE)


# ════════════════════════════════════════════════════════════════
# Slide 7 · DATABASE — 31 tables in 10 domains
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 7, "06 · DATABASE STRUCTURE",
             "Neon Postgres — 31 tables in 10 domains",
             "Single Drizzle schema (db/schema.ts, 591 LOC). Every clinical row carries patient_id ON DELETE CASCADE — GDPR-clean.")

# 10 domain cards in a 5x2 grid
domains = [
    ("01", "Identity & access",       "4", "users · doctor_profiles · patient_profiles · patient_access"),
    ("02", "Clinical structured",     "9", "medications · supplements · surgeries · injuries · clinical_history · risk_assessments · lab_results · vitals_daily · glucose_points"),
    ("03", "Patient artifacts (R2)",  "3", "imaging_studies · writings · documents"),
    ("04", "Self-assessment",         "1", "wheel_of_life_assessments"),
    ("05", "Pipeline (imports)",      "2", "imports · import_files"),
    ("06", "Audit",                   "1", "audit_log"),
    ("07", "Psych architecture",      "3", "psych_dimensions (reference) · psych_items · psych_evidence"),
    ("08", "Mental state & events",   "2", "mood_entries · panic_events"),
    ("09", "Encounters & Rx",         "3", "encounters · prescriptions · taper_history"),
    ("10", "ECG · PGx · life events", "3", "ecg_events · pgx_findings · life_events"),
]

grid_x0 = Inches(0.5)
grid_y0 = Inches(2.1)
card_w  = Inches(2.5)
card_h  = Inches(2.4)
gap_x   = Inches(0.07)
gap_y   = Inches(0.15)

for i, (num, title, n, tables) in enumerate(domains):
    col = i % 5
    row = i // 5
    x = grid_x0 + (card_w + gap_x) * col
    y = grid_y0 + (card_h + gap_y) * row
    add_rect(s, x, y, card_w, card_h, fill=WHITE, stroke=BORDER)
    # gold top stripe
    add_rect(s, x, y, card_w, Inches(0.07), fill=GOLD, stroke=None)
    # number
    add_text(s, x + Inches(0.15), y + Inches(0.18), Inches(0.6), Inches(0.3),
             num, font=FONT_MONO, size=10, bold=True, color=GOLD)
    # count pill (right)
    add_rect(s, x + card_w - Inches(0.6), y + Inches(0.18), Inches(0.45), Inches(0.3),
             fill=NAVY, stroke=None, rounded=True)
    add_text(s, x + card_w - Inches(0.6), y + Inches(0.21), Inches(0.45), Inches(0.25),
             n, font=FONT_MONO, size=10, bold=True, color=WHITE, align="center")
    # title
    add_text(s, x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), Inches(0.5),
             title, font=FONT_DISPLAY, size=12, bold=True, color=NAVY)
    # divider
    add_line(s, x + Inches(0.15), y + Inches(1.0),
             x + card_w - Inches(0.15), y + Inches(1.0), color=BORDER)
    # tables (italic-ish list)
    add_text(s, x + Inches(0.15), y + Inches(1.08), card_w - Inches(0.3), card_h - Inches(1.2),
             tables, font=FONT_MONO, size=8.5, color=BLUE_700, line_spacing=1.35)

# Bottom convention bar
add_rect(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         fill=NAVY, stroke=None)
add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.25),
         "PATIENT_ID FK ON DELETE CASCADE   ·   UUID v4 PKs (bigint for high-volume)   ·   TIMESTAMPTZ EVERYWHERE   ·   9 ENUMS · JSONB FOR FLEX",
         font=FONT_MONO, size=9, bold=True, color=GOLD, anchor="middle")


# ════════════════════════════════════════════════════════════════
# Slide 8 · DATABASE — 3 storage shapes
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 8, "06 · DATABASE STRUCTURE",
             "Three storage shapes — sometimes combined in the same row",
             "Strict columns where you query · JSONB where shape varies · R2 where bytes shouldn't bloat the database.")

shapes_data = [
    {
        "eyebrow": "PATTERN 01 · STRUCTURED",
        "title":   "Strict typed columns",
        "sub":     "One row = one fact. Plottable.",
        "fill":    WHITE, "stroke": NAVY, "stroke_w": 1.4,
        "table_name": "medications",
        "table_color": BLUE_700,
        "columns": [
            ("id",           "uuid · pk"),
            ("patient_id",   "uuid · fk → users.id"),
            ("name",         "text · not null"),
            ("dose",         "text"),
            ("started_at",   "date"),
            ("ended_at",     "date"),
        ],
        "highlight_idx": None,
        "other": "vitals_daily · glucose_points · supplements · surgeries · injuries · lab_results · prescriptions · taper_history · life_events · pgx_findings",
    },
    {
        "eyebrow": "PATTERN 02 · SEMI-STRUCTURED",
        "title":   "Typed envelope + JSONB",
        "sub":     "Schema where it fits, flex where it doesn't.",
        "fill":    GOLD_SOFT, "stroke": GOLD, "stroke_w": 1.6,
        "table_name": "risk_assessments",
        "table_color": GOLD,
        "columns": [
            ("id",          "uuid · pk"),
            ("patient_id",  "uuid · fk"),
            ("kind",        "text · 'AUDIT' · 'PHQ-9'"),
            ("payload",     "jsonb · not null"),
            ("recorded_at", "timestamptz"),
        ],
        "highlight_idx": 3,
        "other": "wheel_of_life_assessments (16 axes in scores) · vitals_daily.extras · panic_events.symptoms · audit_log.metadata · imports.errors",
    },
    {
        "eyebrow": "PATTERN 03 · BLOB-BACKED",
        "title":   "Postgres pointer · R2 bytes",
        "sub":     "Row stays small. Binary content in R2.",
        "fill":    BLUE_50, "stroke": NAVY, "stroke_w": 1.4,
        "table_name": "documents",
        "table_color": BLUE_800,
        "columns": [
            ("id",                "uuid · pk"),
            ("patient_id",        "uuid · fk"),
            ("kind",              "text · 'genetics' · 'lab'"),
            ("original_filename", "text"),
            ("blob_key",          "text → R2 object"),
            ("size_bytes",        "bigint"),
        ],
        "highlight_idx": 4,
        "other": "imaging_studies (blob_prefix → multi-file) · writings · ecg_events · import_files · plus source_blob_key on most clinical tables",
    },
]

col_w = Inches(4.15)
col_h = Inches(4.6)
y0 = Inches(2.15)
x0 = Inches(0.5)
gap = Inches(0.13)

for i, d in enumerate(shapes_data):
    x = x0 + (col_w + gap) * i
    add_rect(s, x, y0, col_w, col_h, fill=d["fill"], stroke=d["stroke"], stroke_w=d["stroke_w"])
    # eyebrow
    add_text(s, x + Inches(0.18), y0 + Inches(0.18), col_w - Inches(0.36), Inches(0.25),
             d["eyebrow"], font=FONT_MONO, size=9, bold=True, color=GOLD)
    # title
    add_text(s, x + Inches(0.18), y0 + Inches(0.45), col_w - Inches(0.36), Inches(0.35),
             d["title"], font=FONT_DISPLAY, size=14, bold=True, color=NAVY)
    # sub
    add_text(s, x + Inches(0.18), y0 + Inches(0.82), col_w - Inches(0.36), Inches(0.25),
             d["sub"], font=FONT_BODY, size=10.5, color=SLATE)

    # table header band
    table_y = y0 + Inches(1.18)
    add_rect(s, x + Inches(0.18), table_y, col_w - Inches(0.36), Inches(0.28),
             fill=d["table_color"], stroke=None)
    add_text(s, x + Inches(0.28), table_y + Inches(0.05), col_w - Inches(0.6), Inches(0.22),
             d["table_name"], font=FONT_MONO, size=10, bold=True, color=WHITE, anchor="middle")
    add_text(s, x + Inches(0.28), table_y + Inches(0.05), col_w - Inches(0.5), Inches(0.22),
             "type", font=FONT_MONO, size=8.5,
             color=RGBColor(0xFF, 0xFF, 0xFF), anchor="middle", align="right")

    # columns
    row_y = table_y + Inches(0.31)
    for ci, (col_name, col_type) in enumerate(d["columns"]):
        is_hl = (d["highlight_idx"] == ci)
        if is_hl:
            row_fill = GOLD_SOFT if i == 1 else (BLUE_50 if i == 2 else WHITE)
            row_stroke = GOLD if i == 1 else (BLUE_500 if i == 2 else BORDER)
            name_color = NAVY
            type_color = NAVY
            font_bold = True
        else:
            row_fill = WHITE
            row_stroke = BORDER
            name_color = SLATE
            type_color = BLUE_700
            font_bold = False
        add_rect(s, x + Inches(0.18), row_y, col_w - Inches(0.36), Inches(0.26),
                 fill=row_fill, stroke=row_stroke,
                 stroke_w=(1.4 if is_hl else 0.6))
        add_text(s, x + Inches(0.28), row_y + Inches(0.04),
                 col_w - Inches(0.6) - Inches(1.2), Inches(0.2),
                 col_name, font=FONT_MONO, size=9.5,
                 bold=font_bold, color=name_color, anchor="middle")
        add_text(s, x + col_w - Inches(1.5), row_y + Inches(0.04), Inches(1.3), Inches(0.2),
                 col_type, font=FONT_MONO, size=9,
                 bold=font_bold, color=type_color,
                 align="right", anchor="middle")
        row_y += Inches(0.28)

    # "OTHER TABLES" footer
    foot_y = y0 + col_h - Inches(1.2)
    add_line(s, x + Inches(0.18), foot_y,
             x + col_w - Inches(0.18), foot_y, color=BORDER)
    add_text(s, x + Inches(0.18), foot_y + Inches(0.08), col_w - Inches(0.36), Inches(0.2),
             "OTHER TABLES OF THIS SHAPE",
             font=FONT_MONO, size=8.5, bold=True, color=STEEL)
    add_text(s, x + Inches(0.18), foot_y + Inches(0.3), col_w - Inches(0.36), Inches(0.85),
             d["other"], font=FONT_MONO, size=8.5, color=BLUE_700, line_spacing=1.3)

# Bottom strip
add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5),
         fill=NAVY, stroke=None)
add_text(s, Inches(0.7), Inches(7.02), Inches(12), Inches(0.4),
         "COMBINED PATTERN · most rows carry all three: typed columns for what you query · optional jsonb for instrument-specific data · optional source_blob_key for audit back to the original PDF",
         font=FONT_MONO, size=10, color=GOLD, anchor="middle", line_spacing=1.25)


# ════════════════════════════════════════════════════════════════
# Slide 9 · AI LAYER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 9, "07 · AI LAYER",
             "Three Claude models, three jobs",
             "Different tasks have different latency / accuracy budgets. One API key, three model identifiers, zero auto-routing.")

models = [
    {"tier": "OPUS 4.7",
     "id": "claude-opus-4-7",
     "purpose": "Patient-facing chat",
     "called_from": "handleChat()  ·  web/_worker.js",
     "rationale": "User-visible reasoning quality matters most. Strict system prompt: answer only from the bundled patient-record.txt — never invent facts.",
     "stripe": NAVY},
    {"tier": "SONNET 4.6",
     "id": "claude-sonnet-4-6",
     "purpose": "Lab PDF transcription · document ingest · dashboard composition",
     "called_from": "lib/ingest.js  ·  lib/dashboard.js",
     "rationale": "Strong vision + structured-output reliability. Used for the heavy lifting that lands rows in Postgres.",
     "stripe": GOLD},
    {"tier": "HAIKU 4.5",
     "id": "claude-haiku-4-5-20251001",
     "purpose": "Document classification · reclassification · doctor-name backfill",
     "called_from": "lib/ingest.js (taxonomy steps)",
     "rationale": "High-volume, low-stakes calls. Latency and cost matter; accuracy is a bonus.",
     "stripe": BLUE_500},
]

card_w = Inches(4.15)
card_h = Inches(4.5)
y0 = Inches(2.2)
x0 = Inches(0.5)
gap = Inches(0.13)

for i, m in enumerate(models):
    x = x0 + (card_w + gap) * i
    add_rect(s, x, y0, card_w, card_h, fill=WHITE, stroke=BORDER)
    # left stripe
    add_rect(s, x, y0, Inches(0.1), card_h, fill=m["stripe"], stroke=None)

    add_text(s, x + Inches(0.3), y0 + Inches(0.2), card_w - Inches(0.5), Inches(0.3),
             m["tier"], font=FONT_MONO, size=11, bold=True, color=GOLD)
    add_text(s, x + Inches(0.3), y0 + Inches(0.55), card_w - Inches(0.5), Inches(0.55),
             "Claude " + m["tier"].title(), font=FONT_DISPLAY, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y0 + Inches(1.1), card_w - Inches(0.5), Inches(0.25),
             m["id"], font=FONT_MONO, size=10, color=BLUE_700)

    add_line(s, x + Inches(0.3), y0 + Inches(1.45),
             x + card_w - Inches(0.3), y0 + Inches(1.45), color=BORDER)

    add_text(s, x + Inches(0.3), y0 + Inches(1.55), card_w - Inches(0.5), Inches(0.22),
             "PURPOSE", font=FONT_MONO, size=8.5, bold=True, color=STEEL)
    add_text(s, x + Inches(0.3), y0 + Inches(1.78), card_w - Inches(0.5), Inches(0.55),
             m["purpose"], font=FONT_BODY, size=12, bold=True, color=NAVY, line_spacing=1.25)

    add_text(s, x + Inches(0.3), y0 + Inches(2.45), card_w - Inches(0.5), Inches(0.22),
             "CALLED FROM", font=FONT_MONO, size=8.5, bold=True, color=STEEL)
    add_text(s, x + Inches(0.3), y0 + Inches(2.68), card_w - Inches(0.5), Inches(0.3),
             m["called_from"], font=FONT_MONO, size=10, color=BLUE_700)

    add_text(s, x + Inches(0.3), y0 + Inches(3.1), card_w - Inches(0.5), Inches(0.22),
             "WHY THIS TIER", font=FONT_MONO, size=8.5, bold=True, color=STEEL)
    add_text(s, x + Inches(0.3), y0 + Inches(3.33), card_w - Inches(0.5), Inches(1.1),
             m["rationale"], font=FONT_BODY, size=11, color=SLATE, line_spacing=1.35)

# Bottom callout
add_rect(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         fill=GOLD_SOFT, stroke=GOLD_LINE)
add_rect(s, Inches(0.5), Inches(7.05), Inches(0.06), Inches(0.35),
         fill=GOLD, stroke=None)
add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.27),
         "ONE KEY · THREE MODELS · NO AUTO-ROUTING.  Each call site pins its own model identifier — choosing the wrong tier for a job is a code change, not a config change.",
         font=FONT_BODY, size=11, bold=False, color=SLATE, anchor="middle")


# ════════════════════════════════════════════════════════════════
# Slide 10 · AUTH & ACCESS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 10, "08 · AUTH & ACCESS",
             "Identity, sessions, role-based access",
             "Custom cookie-session layer (lib/auth.js, Clerk-style API). Every API handler calls authenticate(request, env) before touching data.")

roles = [
    {"name": "ADMIN",
     "title": "Platform administrator",
     "scope": "Full read / write across every patient",
     "bullets": [
         "Creates other users (doctors + patients)",
         "Deletes patients (cascades through Postgres + R2)",
         "Access: /api/admin/*",
         "Only role that can issue right-to-erasure",
     ],
     "color": NAVY},
    {"name": "DOCTOR",
     "title": "Clinical user",
     "scope": "Read access to patients linked via patient_access",
     "bullets": [
         "Can write dashboard sections + clinical notes",
         "No create-user / delete-user rights",
         "Access: /api/patient-* for linked patients only",
         "Auditable — every action lands in audit_log",
     ],
     "color": GOLD},
    {"name": "PATIENT",
     "title": "Self-access user",
     "scope": "Read access to their own record only",
     "bullets": [
         "Can upload documents via the ingest path",
         "Can wipe their own data (right to erasure)",
         "Access: chat + their own static pages",
         "Bilingual UI · EN / PT toggle in nav",
     ],
     "color": BLUE_500},
]

card_w = Inches(4.15)
card_h = Inches(4.4)
y0 = Inches(2.2)
x0 = Inches(0.5)
gap = Inches(0.13)

for i, r in enumerate(roles):
    x = x0 + (card_w + gap) * i
    add_rect(s, x, y0, card_w, card_h, fill=WHITE, stroke=BORDER)
    # top stripe
    add_rect(s, x, y0, card_w, Inches(0.55), fill=r["color"], stroke=None)
    add_text(s, x + Inches(0.3), y0 + Inches(0.16), card_w - Inches(0.6), Inches(0.3),
             r["name"], font=FONT_MONO, size=14, bold=True, color=WHITE)

    add_text(s, x + Inches(0.3), y0 + Inches(0.78), card_w - Inches(0.6), Inches(0.4),
             r["title"], font=FONT_DISPLAY, size=18, bold=True, color=NAVY)

    add_text(s, x + Inches(0.3), y0 + Inches(1.3), card_w - Inches(0.6), Inches(0.7),
             r["scope"], font=FONT_BODY, size=12, italic=True, color=SLATE, line_spacing=1.35)

    add_line(s, x + Inches(0.3), y0 + Inches(2.15),
             x + card_w - Inches(0.3), y0 + Inches(2.15), color=BORDER)

    by = y0 + Inches(2.3)
    for b in r["bullets"]:
        # bullet dot
        add_rect(s, x + Inches(0.35), by + Inches(0.08), Inches(0.08), Inches(0.08),
                 fill=GOLD, stroke=None, rounded=True)
        add_text(s, x + Inches(0.55), by, card_w - Inches(0.85), Inches(0.45),
                 b, font=FONT_BODY, size=11, color=SLATE, line_spacing=1.3)
        by += Inches(0.5)

# Bottom flow callout
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         fill=NAVY, stroke=None)
add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.32),
         "REQUEST → AUTHENTICATE(REQUEST, ENV) → ROLE CHECK → PATIENT_ACCESS CHECK → HANDLER → AUDIT_LOG ROW",
         font=FONT_MONO, size=10, bold=True, color=GOLD, anchor="middle")


# ════════════════════════════════════════════════════════════════
# Slide 11 · DEPLOYMENT FLOW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 11, "09 · DEPLOYMENT",
             "Source vs. deploy — two separate operations",
             "GitHub is the source of truth. It does NOT deploy. Live URL only updates when an engineer runs the Wrangler CLI against web/.")

# Engineer box
add_rect(s, Inches(0.5), Inches(3.4), Inches(2.4), Inches(1.4),
         fill=WHITE, stroke=NAVY)
add_text(s, Inches(0.5), Inches(3.55), Inches(2.4), Inches(0.35),
         "Engineer · laptop", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(0.5), Inches(3.95), Inches(2.4), Inches(0.75),
         "edits source +\nruns ETL locally\n(python3 bin/extract.py)",
         font=FONT_BODY, size=11, color=SLATE, align="center", line_spacing=1.3)

# Branch 1: GitHub (top)
add_arrow(s, Inches(2.95), Inches(3.7), Inches(3.7), Inches(2.85), color=SLATE)
add_text(s, Inches(3.0), Inches(3.25), Inches(0.9), Inches(0.25),
         "git push", font=FONT_MONO, size=10, bold=True, color=SLATE, align="center")

add_rect(s, Inches(3.75), Inches(2.4), Inches(3.2), Inches(1.4),
         fill=WHITE, stroke=STEEL, dash=True)
add_text(s, Inches(3.75), Inches(2.55), Inches(3.2), Inches(0.35),
         "GitHub", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(3.75), Inches(2.95), Inches(3.2), Inches(0.3),
         "joaocreste/health", font=FONT_MONO, size=10, color=BLUE_700, align="center")
add_text(s, Inches(3.75), Inches(3.3), Inches(3.2), Inches(0.3),
         "source-of-truth only",
         font=FONT_BODY, size=11, italic=True, color=SLATE, align="center")
add_text(s, Inches(3.75), Inches(3.55), Inches(3.2), Inches(0.25),
         "DOES NOT DEPLOY", font=FONT_MONO, size=10, bold=True, color=GOLD, align="center")

# Branch 2: Wrangler → CF Pages → User (bottom)
add_arrow(s, Inches(2.95), Inches(4.5), Inches(3.7), Inches(5.35), color=GOLD, weight=1.8)
add_text(s, Inches(2.85), Inches(4.95), Inches(1.7), Inches(0.25),
         "wrangler pages deploy", font=FONT_MONO, size=10, bold=True, color=GOLD, align="center")

# Wrangler
add_rect(s, Inches(3.75), Inches(5.0), Inches(3.2), Inches(1.4),
         fill=GOLD_SOFT, stroke=GOLD, stroke_w=1.6)
add_text(s, Inches(3.75), Inches(5.15), Inches(3.2), Inches(0.35),
         "Wrangler CLI", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(3.75), Inches(5.55), Inches(3.2), Inches(0.3),
         "uploads web/ folder", font=FONT_BODY, size=11, color=SLATE, align="center")
add_text(s, Inches(3.75), Inches(5.85), Inches(3.2), Inches(0.3),
         "--commit-dirty=true", font=FONT_MONO, size=10, color=BLUE_700, align="center")

# → CF Pages
add_arrow(s, Inches(7.0), Inches(5.7), Inches(7.4), Inches(5.7), color=GOLD, weight=1.8)

add_rect(s, Inches(7.45), Inches(5.0), Inches(3.0), Inches(1.4),
         fill=NAVY, stroke=NAVY)
add_text(s, Inches(7.45), Inches(5.15), Inches(3.0), Inches(0.35),
         "Cloudflare Pages", font=FONT_DISPLAY, size=14, bold=True, color=WHITE, align="center")
add_text(s, Inches(7.45), Inches(5.55), Inches(3.0), Inches(0.3),
         "jc-advisory-health.pages.dev", font=FONT_MONO, size=10, color=GOLD, align="center")
add_text(s, Inches(7.45), Inches(5.85), Inches(3.0), Inches(0.3),
         "unique preview URL per deploy",
         font=FONT_BODY, size=10, color=RGBColor(0xCF, 0xD7, 0xE0), align="center")

# → User
add_arrow(s, Inches(10.5), Inches(5.7), Inches(10.9), Inches(5.7), color=SLATE)

add_rect(s, Inches(10.95), Inches(5.0), Inches(1.9), Inches(1.4),
         fill=WHITE, stroke=NAVY)
add_text(s, Inches(10.95), Inches(5.15), Inches(1.9), Inches(0.35),
         "User · live", font=FONT_DISPLAY, size=14, bold=True, color=NAVY, align="center")
add_text(s, Inches(10.95), Inches(5.55), Inches(1.9), Inches(0.3),
         "global CDN", font=FONT_BODY, size=11, color=SLATE, align="center")
add_text(s, Inches(10.95), Inches(5.85), Inches(1.9), Inches(0.3),
         "≈seconds after deploy", font=FONT_BODY, size=10, color=SLATE, align="center")

# Daily-flow strip at bottom
add_rect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55),
         fill=SURFACE, stroke=BORDER)
add_text(s, Inches(0.7), Inches(6.92), Inches(12), Inches(0.22),
         "DAILY FLOW (always both, in either order):",
         font=FONT_MONO, size=9, bold=True, color=STEEL)
add_text(s, Inches(0.7), Inches(7.15), Inches(12), Inches(0.25),
         "git add … → git commit (ASCII arrows only) → git push   ·   npx wrangler pages deploy web --commit-dirty=true   ·   curl preview URL to verify",
         font=FONT_MONO, size=10, color=BLUE_700)


# ════════════════════════════════════════════════════════════════
# Slide 12 · SECURITY POSTURE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 12, "10 · SECURITY & COMPLIANCE",
             "Where we are, where we need to get to",
             "Architecture is shaped for HIPAA / LGPD / GDPR. No formal certification yet — sized for <10 patients pre-certification.")

posture_cols = [
    {"title": "IN PLACE TODAY",
     "color": BLUE_500,
     "fill": BLUE_50,
     "bullets": [
         "EU-only Postgres region (Frankfurt)",
         "Row-level patient ownership · ON DELETE CASCADE",
         "R2 per-patient key prefixes",
         "Audit log of every mutation",
         "Role-based API gating on every Worker handler",
         "HTTPS-only · Cloudflare-managed TLS",
         "Right-to-erasure cascades Postgres + R2",
     ]},
    {"title": "DEFERRED UNTIL REAL PATIENTS",
     "color": GOLD,
     "fill": GOLD_SOFT,
     "bullets": [
         "HIPAA Business Associate Agreement (Cloudflare)",
         "HIPAA BAA with Anthropic (Enterprise plan)",
         "Neon Enterprise (encryption-at-rest options)",
         "SOC 2 / ISO 27001 attestations",
         "Formal DPIA + LGPD Article 50 documentation",
         "Penetration test on the Worker surface",
     ]},
    {"title": "COMPENSATING CONTROLS",
     "color": NAVY,
     "fill": SURFACE,
     "bullets": [
         "Patient writings stay local until ingested",
         "866 MB Apple Health XML is gitignored",
         "Wearable exports kept off-repo for non-Joao patients",
         "Anthropic system prompt forbids inventing facts",
     ]},
]

col_w = Inches(4.15)
col_h = Inches(4.4)
y0 = Inches(2.2)
x0 = Inches(0.5)
gap = Inches(0.13)

for i, c in enumerate(posture_cols):
    x = x0 + (col_w + gap) * i
    add_rect(s, x, y0, col_w, col_h, fill=c["fill"], stroke=c["color"], stroke_w=1.3)
    add_rect(s, x, y0, col_w, Inches(0.5), fill=c["color"], stroke=None)
    add_text(s, x + Inches(0.25), y0 + Inches(0.13), col_w - Inches(0.5), Inches(0.3),
             c["title"], font=FONT_MONO, size=11, bold=True, color=WHITE)
    by = y0 + Inches(0.75)
    for b in c["bullets"]:
        add_rect(s, x + Inches(0.32), by + Inches(0.13), Inches(0.06), Inches(0.06),
                 fill=GOLD, stroke=None, rounded=True)
        add_text(s, x + Inches(0.5), by, col_w - Inches(0.7), Inches(0.5),
                 b, font=FONT_BODY, size=11, color=SLATE, line_spacing=1.3)
        by += Inches(0.5)

# Posture summary callout
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         fill=NAVY, stroke=None)
add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.32),
         "CURRENT POSTURE · single-engineer · single live patient (plus 3 bespoke demos) · pre-revenue · pre-certification · upgrades to Enterprise tiers happen once real patients sign on",
         font=FONT_MONO, size=10, color=GOLD, anchor="middle")


# ─── Write file ──────────────────────────────────────────────────
OUT = Path(__file__).resolve().parent.parent / "infrastructure.pptx"
prs.save(OUT)
print(f"Wrote {OUT} ({OUT.stat().st_size:,} bytes)")
