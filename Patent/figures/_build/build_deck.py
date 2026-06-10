# -*- coding: utf-8 -*-
"""Assemble the Lumen Health figure set into a 16:9 branded PPTX deck.
Title slide + EN section (9 figs) + PT section (9 figs). Figures are the
300dpi PNGs already on disk; each is fit to its own slide on a light ground."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)                       # Patent/figures
LOGO = os.path.join(HERE, "_logo.png")
OUT  = os.path.join(ROOT, "Lumen-Health-Modus-Operandi.pptx")

# brand (dark title ground + light figure ground)
NAVY   = RGBColor(0x0A, 0x14, 0x28)
SLATE  = RGBColor(0x13, 0x1F, 0x37)
GOLD_D = RGBColor(0xF4, 0xB9, 0x42)
GOLD_L = RGBColor(0xB8, 0x86, 0x0B)
WHITEISH = RGBColor(0xF0, 0xF4, 0xF8)
LIGHTBG = RGBColor(0xF7, 0xF8, 0xFA)
PETROL = RGBColor(0x24, 0x4E, 0x6E)
INK    = RGBColor(0x0D, 0x1B, 0x33)
MUTED  = RGBColor(0x6F, 0x7E, 0x9A)

EW, EH = Inches(13.333), Inches(7.5)

FIGS = [
    ("1", "end-to-end",            "fluxo-geral",          "End-to-end method flow",          "Fluxo geral ponta a ponta"),
    ("2", "patient-portal",        "paciente-portal",      "Patient to portal",               "Paciente ao portal"),
    ("3", "etl-cloud",             "etl-nuvem",            "ETL pipeline",                    "Pipeline ETL"),
    ("4", "classification-prompts","classificacao-prompts","Classification & prompt registry","Classificacao e registro de prompts"),
    ("5", "import-state-machine",  "estados-importacao",   "Import job state machine",        "Maquina de estados da importacao"),
    ("6", "three-dimensions",      "tres-dimensoes",       "Three-dimension organization",    "Organizacao em tres dimensoes"),
    ("7", "record-two-views",      "registro-duas-visoes", "One record, two views",           "Um registro, duas visoes"),
    ("8", "per-patient-ai-agent",  "agente-ia-paciente",   "Per-patient AI agent",            "Agente de IA por paciente"),
    ("9", "governance-erasure",    "governanca-exclusao",  "Governance & erasure",            "Governanca e exclusao"),
]

prs = Presentation()
prs.slide_width = EW
prs.slide_height = EH
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold, font, spacing) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        f = r.font
        f.size = Pt(size); f.bold = bold; f.name = font
        f.color.rgb = color
    return tb

def fit_image(slide, path, top_margin=Inches(0.35), bot_margin=Inches(0.45)):
    iw, ih = Image.open(path).size
    avail_h = EH - top_margin - bot_margin
    h = avail_h
    w = Emu(int(int(h) * iw / ih))
    if w > EW - Inches(0.6):                # cap width for landscape
        w = EW - Inches(0.6)
        h = Emu(int(int(w) * ih / iw))
    x = Emu(int((int(EW) - int(w)) / 2))
    y = Emu(int((int(EH) - int(h)) / 2) - int(Inches(0.06)))
    slide.shapes.add_picture(path, x, y, w, h)

def footer(slide, idx, dark=False):
    col = MUTED if not dark else RGBColor(0x8A, 0x9B, 0xB5)
    textbox(slide, Inches(0.45), EH - Inches(0.42), Inches(8), Inches(0.3),
            [("LUMEN HEALTH  ·  MODUS OPERANDI  ·  CONFIDENTIAL", 8, col, False,
              "Consolas", 0)], align=PP_ALIGN.LEFT)
    textbox(slide, EW - Inches(1.4), EH - Inches(0.42), Inches(0.95), Inches(0.3),
            [(str(idx), 8, col, False, "Consolas", 0)], align=PP_ALIGN.RIGHT)

# ---- title slide ----
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
s.shapes.add_picture(LOGO, Inches(5.9), Inches(1.5), Inches(1.55), Inches(1.55))
textbox(s, Inches(1), Inches(3.15), Inches(11.333), Inches(1.0),
        [("Lumen Health", 46, RGBColor(0xFF,0xFF,0xFF), False, "Raleway Light", 0)],
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.15), Inches(11.333), Inches(0.5),
        [("MODUS OPERANDI  ·  FIGURE SET", 13, GOLD_D, False, "Consolas", 0)],
        align=PP_ALIGN.CENTER)
textbox(s, Inches(2.5), Inches(5.0), Inches(8.333), Inches(0.8),
        [("Data-agnostic ingestion, ETL pipeline and per-patient AI insight layer.",
          14, RGBColor(0xB5,0xC2,0xD6), False, "Calibri", 0),
         ("9 figures  ·  English + Portugues  ·  branded light theme", 11,
          RGBColor(0x8A,0x9B,0xB5), False, "Consolas", 0)],
        align=PP_ALIGN.CENTER)
footer(s, 1, dark=True)

def divider(title, sub, idx):
    s = prs.slides.add_slide(BLANK); bg(s, SLATE)
    textbox(s, Inches(1.2), Inches(3.0), Inches(11), Inches(1.1),
            [(title, 38, RGBColor(0xFF,0xFF,0xFF), False, "Raleway Light", 0)],
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(1.25), Inches(4.1), Inches(11), Inches(0.5),
            [(sub, 13, GOLD_D, False, "Consolas", 0)], align=PP_ALIGN.LEFT)
    footer(s, idx, dark=True)

def figure_slide(png, title, fig_no, idx):
    s = prs.slides.add_slide(BLANK); bg(s, LIGHTBG)
    fit_image(s, png)
    footer(s, idx)

n = 1
# EN section
n += 1; divider("English figures", "FIG. 1 - 9  ·  for PCT / EPO reading", n)
for no, en_slug, pt_slug, en_t, pt_t in FIGS:
    n += 1
    figure_slide(os.path.join(ROOT, "en", "png", f"fig-{no}-{en_slug}.png"),
                 en_t, no, n)
# PT section
n += 1; divider("Figuras em portugues", "FIG. 1 - 9  ·  para leitura INPI", n)
for no, en_slug, pt_slug, en_t, pt_t in FIGS:
    n += 1
    figure_slide(os.path.join(ROOT, "pt", "png", f"fig-{no}-{pt_slug}.png"),
                 pt_t, no, n)

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
